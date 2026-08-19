from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


OUT_PATH = r"C:\Users\RIPOLN\OneDrive - Pfizer\Documents\Vyndamax\Dx headroom\ATTR_CM_Future_NP_Deck.pptx"

BLUE = RGBColor(0, 94, 184)
DARK = RGBColor(12, 29, 51)
LIGHT = RGBColor(245, 248, 252)
ACCENT = RGBColor(0, 164, 224)
TEXT = RGBColor(32, 40, 48)
MUTED = RGBColor(92, 102, 112)
GREEN = RGBColor(0, 128, 96)
RED = RGBColor(184, 32, 48)
WHITE = RGBColor(255, 255, 255)


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None, accent=False):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = BLUE if accent else DARK
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.92), Inches(11.2), Inches(0.35))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(11)
        run2.font.color.rgb = MUTED

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.35), Inches(11.8), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()


def add_footer(slide, text="Pfizer"):
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(1.6), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = BLUE


def add_bullet_list(slide, bullets, left, top, width, height, font_size=18, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        p.font.size = Pt(font_size)
        p.font.color.rgb = color


def add_metric_box(slide, title, value, subtitle, left, top, width, height):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = LIGHT
    sh.line.color.rgb = BLUE
    sh.line.width = Pt(1.2)

    tx = sh.text_frame
    tx.word_wrap = True
    p1 = tx.paragraphs[0]
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    p1.font.bold = True
    p1.font.size = Pt(12)
    p1.font.color.rgb = MUTED

    p2 = tx.add_paragraph()
    p2.text = value
    p2.alignment = PP_ALIGN.CENTER
    p2.font.bold = True
    p2.font.size = Pt(23)
    p2.font.color.rgb = BLUE

    p3 = tx.add_paragraph()
    p3.text = subtitle
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(10)
    p3.font.color.rgb = MUTED


def add_flow_box(slide, title, left, top, width, height, fill_color=LIGHT, text_color=TEXT):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    sh.line.color.rgb = BLUE
    sh.line.width = Pt(1.5)
    tx = sh.text_frame
    tx.word_wrap = True
    p = tx.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = text_color


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = BLUE
    line.line.width = Pt(2)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "Identifying Future ATTR-CM New Patients", "Disease detection strategy using claim-driven patient journeys and predictive risk scoring", accent=True)

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.9), Inches(5.7), Inches(2.4))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Objective: identify patients with a clinically meaningful pre-diagnosis trajectory that may be on the ATTR-CM pathway before formal diagnosis."
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT
    p.space_after = Pt(8)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(1.9), Inches(4.6), Inches(2.1))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = BLUE
    card.line.width = Pt(1.5)
    tx = card.text_frame
    tx.word_wrap = True
    p1 = tx.paragraphs[0]
    p1.text = "Key question"
    p1.font.bold = True; p1.font.size = Pt(14); p1.font.color.rgb = BLUE
    p2 = tx.add_paragraph()
    p2.text = "Can we detect high-risk, undiagnosed patients earlier by learning from longitudinal claim and procedure patterns?"
    p2.font.size = Pt(16); p2.font.color.rgb = TEXT

    add_metric_box(slide, "Reduced candidate cohort", "88.9M", "Clinically filtered", Inches(0.75), Inches(4.8), Inches(2.8), Inches(1.3))
    add_metric_box(slide, "Balanced training", "489k", "Positive + sampled negatives", Inches(3.9), Inches(4.8), Inches(2.9), Inches(1.3))
    add_metric_box(slide, "Baseline ROC AUC", "0.90", "Strong separation", Inches(7.1), Inches(4.8), Inches(2.8), Inches(1.3))
    add_metric_box(slide, "Top-risk use case", "Future NP IDs", "Earlier clinical review", Inches(10.3), Inches(4.8), Inches(2.2), Inches(1.3))
    add_footer(slide)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "Why this matters", "The problem is not just diagnosis coding â€” it is the pattern that precedes diagnosis")
    add_bullet_list(slide, [
        "ATTR-CM is often preceded by a sequence of clinically meaningful signals, not a single trigger event.",
        "Traditional patient-level summaries are too aggregated to reconstruct the true clinical timeline.",
        "The opportunity is to identify future new patients (NPs) by recognizing early red-flag and procedure patterns.",
        "A longitudinal journey model provides a better view of disease progression than a static patient snapshot."
    ], Inches(0.8), Inches(1.9), Inches(6.5), Inches(3.4), font_size=20)

    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.9), Inches(4.5), Inches(3.8))
    right_box.fill.solid(); right_box.fill.fore_color.rgb = LIGHT
    right_box.line.color.rgb = BLUE
    right_box.line.width = Pt(1.5)
    tx = right_box.text_frame
    tx.word_wrap = True
    p = tx.paragraphs[0]
    p.text = "Clinical sequence of interest"
    p.font.bold = True; p.font.size = Pt(15); p.font.color.rgb = BLUE
    for i, item in enumerate([
        "First claim",
        "AFib / HF / HFpEF signals",
        "Red-flag clustering",
        "ATTR-related procedures",
        "First ATTR-CM diagnosis"
    ]):
        p2 = tx.add_paragraph()
        p2.text = item
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT
        p2.bullet = True
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "How the patient journey works", "Claim and procedure data are transformed into a longitudinal clinical timeline")

    add_flow_box(slide, "1. Raw claims\n& procedures", Inches(0.7), Inches(2.1), Inches(2.2), Inches(1.3), fill_color=LIGHT)
    add_flow_box(slide, "2. Event extraction\nAFib, HF, HFpEF, red flags", Inches(3.2), Inches(2.1), Inches(2.7), Inches(1.3), fill_color=LIGHT)
    add_flow_box(slide, "3. Procedure signals\nPYP, IHC, MRI, ECG, biopsy, etc.", Inches(6.2), Inches(2.1), Inches(2.8), Inches(1.3), fill_color=LIGHT)
    add_flow_box(slide, "4. Milestone table\nfirst dates & lead times", Inches(9.4), Inches(2.1), Inches(2.8), Inches(1.3), fill_color=LIGHT)

    add_bullet_list(slide, [
        "We use claim-level diagnosis tables to create patient-by-event timelines.",
        "We map procedure claims to clinically meaningful ATTR-related event families.",
        "The milestone table summarizes timing and cumulative risk signals at the patient level.",
        "This creates a clean foundation for risk scoring and future NP detection."
    ], Inches(1.0), Inches(4.2), Inches(11.0), Inches(2.0), font_size=18)
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "Cohort and model results", "The model is trained on a clinically enriched, high-risk population")

    add_metric_box(slide, "Raw source population", "335M+", "All patients in source claims", Inches(0.8), Inches(2.0), Inches(2.6), Inches(1.5))
    add_metric_box(slide, "Reduced candidate cohort", "88.9M", "HF / AFib / red flags / procedures", Inches(3.8), Inches(2.0), Inches(2.9), Inches(1.5))
    add_metric_box(slide, "Balanced training set", "489k", "81.6k positives, 408k negatives", Inches(7.2), Inches(2.0), Inches(2.9), Inches(1.5))
    add_metric_box(slide, "Model performance", "ROC AUC 0.90", "Avg precision 0.78", Inches(10.5), Inches(2.0), Inches(2.2), Inches(1.5))

    add_bullet_list(slide, [
        "We intentionally narrowed the cohort to clinically relevant patients rather than modeling the full claims universe.",
        "This supports better signal quality and a more interpretable early-risk strategy.",
        "The first-pass model already shows strong discrimination between diagnosed and similar undiagnosed patients."
    ], Inches(1.0), Inches(4.4), Inches(11.0), Inches(1.8), font_size=18)
    add_footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "How this helps identify future NPs", "Turn patient risk trajectories into actionable opportunity lists")

    add_bullet_list(slide, [
        "Step 1: score all candidate patients using the learned pre-diagnosis signal pattern.",
        "Step 2: rank undiagnosed patients by risk and examine top deciles for likely early ATTR-CM pathways.",
        "Step 3: prioritize patients with AFib + HF/HFpEF + procedure clustering for clinical review.",
        "Step 4: monitor the same journey over time to detect new patients as their risk signal accumulates.",
        "Step 5: enrich field force and care strategy by focusing on the patients most likely to convert to ATTR-CM diagnosis."
    ], Inches(0.9), Inches(1.9), Inches(7.0), Inches(3.6), font_size=19)

    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(2.0), Inches(3.8), Inches(2.8))
    callout.fill.solid(); callout.fill.fore_color.rgb = LIGHT
    callout.line.color.rgb = GREEN
    callout.line.width = Pt(1.5)
    tx = callout.text_frame
    tx.word_wrap = True
    p = tx.paragraphs[0]
    p.text = "Future NP use case"
    p.font.bold = True; p.font.size = Pt(15); p.font.color.rgb = GREEN
    p2 = tx.add_paragraph()
    p2.text = "Identify patients with rising red-flag burden before formal ATTR-CM diagnosis and direct intervention or deeper review."
    p2.font.size = Pt(18); p2.font.color.rgb = DARK

    add_footer(slide)

    prs.save(OUT_PATH)
    print(f"Saved deck to: {OUT_PATH}")


if __name__ == "__main__":
    build_deck()

